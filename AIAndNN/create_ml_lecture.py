"""
Script to generate a PowerPoint presentation on Machine Learning
for undergraduate students (2-hour lecture)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ml_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set slide width and height (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Helper function to add title slide
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
    
    # Helper function to add content slide
    def add_content_slide(title, bullet_points, visual_note=""):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = "Calibri"
        
        # Add visual note if provided
        if visual_note:
            p = tf.add_paragraph()
            p.text = f"\n[Visual Suggestion: {visual_note}]"
            p.font.size = Pt(14)
            p.font.italic = True
            p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Helper function to add two-column slide
    def add_two_column_slide(title, left_points, right_points, visual_note=""):
        # Use blank layout for more control
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = title
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.name = "Calibri"
        
        # Add left text box
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        for i, point in enumerate(left_points):
            if i == 0:
                p = left_tf.paragraphs[0]
            else:
                p = left_tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = "Calibri"
        
        # Add right text box
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        
        for i, point in enumerate(right_points):
            if i == 0:
                p = right_tf.paragraphs[0]
            else:
                p = right_tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = "Calibri"
    
    # ============================================
    # TITLE SLIDE
    # ============================================
    add_title_slide(
        "Introduction to Machine Learning\n& Feature Extraction",
        "A Beginner's Guide for Undergraduate Students\n2-Hour Lecture"
    )
    
    # ============================================
    # PART 1: INTRODUCTION TO MACHINE LEARNING
    # ============================================
    
    add_content_slide(
        "Part 1: Introduction to Machine Learning",
        [
            "Duration: 60 minutes",
            "Topics:",
            "• What is Machine Learning?",
            "• Why do we need it?",
            "• Types of ML",
            "• Real-world applications",
            "• ML workflow"
        ],
        "Icon: Brain or computer with learning symbol"
    )
    
    # Slide 2: What is Machine Learning?
    add_content_slide(
        "What is Machine Learning?",
        [
            "Simple Definition:",
            "• Machine Learning is teaching computers to learn from data",
            "• Instead of programming every rule, we let the computer find patterns",
            "",
            "Think of it like this:",
            "• Traditional: You teach a child every single rule",
            "• ML: You show examples and let them figure out the pattern",
            "",
            "Example:",
            "• Show 1000 cat photos → Computer learns what a cat looks like"
        ],
        "Diagram: Traditional programming vs ML (side by side)"
    )
    
    # Slide 3: Why do we need ML?
    add_content_slide(
        "Why do we need Machine Learning?",
        [
            "Problems too complex for traditional programming:",
            "",
            "1. Pattern Recognition",
            "   • Recognizing faces in photos",
            "   • Understanding speech",
            "",
            "2. Too Many Rules",
            "   • Email spam: Rules change constantly",
            "   • Stock market: Millions of factors",
            "",
            "3. Personalization",
            "   • Netflix recommendations",
            "   • YouTube video suggestions",
            "",
            "4. Automation",
            "   • Self-driving cars",
            "   • Medical diagnosis assistance"
        ],
        "Icons: Face recognition, email, recommendation, car"
    )
    
    # Slide 4: Traditional Programming vs ML
    add_two_column_slide(
        "Traditional Programming vs Machine Learning",
        [
            "Traditional Programming:",
            "",
            "Input + Rules → Output",
            "",
            "Example:",
            "• Input: Number",
            "• Rule: If number > 0, output 'Positive'",
            "• Output: 'Positive' or 'Negative'",
            "",
            "Limitations:",
            "• Need to write all rules",
            "• Can't handle complex patterns",
            "• Hard to update"
        ],
        [
            "Machine Learning:",
            "",
            "Input + Output → Rules",
            "",
            "Example:",
            "• Input: 1000 cat photos",
            "• Output: Labels (cat/not cat)",
            "• ML finds the rules automatically",
            "",
            "Advantages:",
            "• Learns from examples",
            "• Handles complex patterns",
            "• Improves with more data"
        ],
        "Flowchart: Two parallel processes"
    )
    
    # Slide 5: Types of Machine Learning
    add_content_slide(
        "Types of Machine Learning",
        [
            "Three main categories:",
            "",
            "1. Supervised Learning",
            "   • Learning with labeled examples",
            "   • Like learning with answer key",
            "",
            "2. Unsupervised Learning",
            "   • Finding patterns in unlabeled data",
            "   • Like exploring without a map",
            "",
            "3. Reinforcement Learning",
            "   • Learning through trial and error",
            "   • Like learning to ride a bike"
        ],
        "Three icons/boxes representing each type"
    )
    
    # Slide 6: Supervised Learning
    add_content_slide(
        "Supervised Learning",
        [
            "Definition:",
            "• We provide labeled training data",
            "• Model learns the relationship",
            "• Then predicts on new data",
            "",
            "Examples:",
            "• Email: Spam (1) or Not Spam (0)",
            "• Images: Cat, Dog, or Bird",
            "• Grades: Predict student performance",
            "",
            "Two types:",
            "• Classification: Categories (spam/not spam)",
            "• Regression: Numbers (house price, temperature)"
        ],
        "Diagram: Labeled data → Model → Prediction"
    )
    
    # Slide 7: Unsupervised Learning
    add_content_slide(
        "Unsupervised Learning",
        [
            "Definition:",
            "• No labels provided",
            "• Find hidden patterns in data",
            "• Discover structure",
            "",
            "Examples:",
            "• Customer segmentation",
            "   → Group similar customers",
            "• Anomaly detection",
            "   → Find unusual transactions",
            "• Image compression",
            "   → Reduce file size while keeping quality",
            "",
            "Common techniques:",
            "• Clustering: Group similar items",
            "• Dimensionality reduction: Simplify data"
        ],
        "Diagram: Unlabeled data → Patterns → Groups"
    )
    
    # Slide 8: Reinforcement Learning
    add_content_slide(
        "Reinforcement Learning",
        [
            "Definition:",
            "• Agent learns by interacting with environment",
            "• Gets rewards for good actions",
            "• Gets penalties for bad actions",
            "",
            "Real-world examples:",
            "• Game playing (Chess, Go)",
            "   → Learn winning strategies",
            "• Robot navigation",
            "   → Learn to avoid obstacles",
            "• Recommendation systems",
            "   → Learn what users like",
            "",
            "Key concept:",
            "• Trial and error with feedback"
        ],
        "Diagram: Agent → Action → Environment → Reward"
    )
    
    # Slide 9: Real-world Examples - Email Spam
    add_content_slide(
        "Real-world Example: Email Spam Detection",
        [
            "Problem:",
            "• Millions of emails daily",
            "• Spam changes constantly",
            "• Can't write rules for everything",
            "",
            "ML Solution:",
            "• Train on thousands of emails",
            "• Learn patterns:",
            "  - Words like 'FREE', 'WIN', 'URGENT'",
            "  - Suspicious sender addresses",
            "  - Unusual formatting",
            "• Classify new emails automatically",
            "",
            "Result:",
            "• Gmail filters 99.9% of spam automatically"
        ],
        "Icons: Email inbox, spam folder, checkmark"
    )
    
    # Slide 10: Real-world Examples - Face Recognition
    add_content_slide(
        "Real-world Example: Face Recognition",
        [
            "Problem:",
            "• Identify people in photos",
            "• Used in:",
            "  - Phone unlock",
            "  - Security systems",
            "  - Social media tagging",
            "",
            "ML Solution:",
            "• Train on millions of face images",
            "• Learn facial features:",
            "  - Eye shape, nose, mouth position",
            "  - Face structure patterns",
            "• Recognize faces in new photos",
            "",
            "How it works:",
            "• Extract facial features",
            "• Compare with known faces",
            "• Match if similar enough"
        ],
        "Diagram: Face → Features → Comparison → Match"
    )
    
    # Slide 11: Real-world Examples - Recommendation Systems
    add_content_slide(
        "Real-world Example: Recommendation Systems",
        [
            "Netflix & YouTube:",
            "• Problem: Too many videos to choose from",
            "• Solution: Recommend what you might like",
            "",
            "How it works:",
            "• Track what you watch",
            "• Find users with similar tastes",
            "• Recommend what they liked",
            "",
            "Example:",
            "• You watched: Action movies",
            "• Similar users also watched: Sci-fi",
            "• Recommendation: Try 'Interstellar'",
            "",
            "Result:",
            "• 80% of Netflix views come from recommendations"
        ],
        "Diagram: User preferences → Similar users → Recommendations"
    )
    
    # Slide 12: ML Workflow
    add_content_slide(
        "Machine Learning Workflow",
        [
            "Step-by-step process:",
            "",
            "1. Data Collection",
            "   • Gather examples (photos, text, numbers)",
            "",
            "2. Data Preprocessing",
            "   • Clean and organize data",
            "   • Remove errors and duplicates",
            "",
            "3. Feature Extraction",
            "   • Convert raw data into features",
            "   • (We'll learn this in Part 2!)",
            "",
            "4. Model Training",
            "   • Teach the model using data",
            "",
            "5. Evaluation",
            "   • Test on new data",
            "   • Measure accuracy",
            "",
            "6. Prediction",
            "   • Use model for real predictions"
        ],
        "Flowchart: Data → Preprocess → Features → Train → Evaluate → Predict"
    )
    
    # Slide 13: Common ML Applications
    add_content_slide(
        "Common Machine Learning Applications",
        [
            "In Daily Life:",
            "• Virtual assistants (Siri, Alexa)",
            "• Social media feeds (Facebook, Instagram)",
            "• Online shopping (Amazon recommendations)",
            "• Navigation (Google Maps route optimization)",
            "",
            "In Industry:",
            "• Healthcare: Disease diagnosis",
            "• Finance: Fraud detection",
            "• Transportation: Self-driving cars",
            "• Education: Personalized learning",
            "",
            "In Research:",
            "• Drug discovery",
            "• Climate prediction",
            "• Language translation"
        ],
        "Grid of icons representing different applications"
    )
    
    # Slide 14: Part 1 Recap
    add_content_slide(
        "Part 1: Key Takeaways",
        [
            "What we learned:",
            "",
            "✓ Machine Learning = Learning from data",
            "✓ Three types: Supervised, Unsupervised, Reinforcement",
            "✓ Used everywhere: Email, photos, recommendations",
            "✓ Workflow: Data → Features → Model → Prediction",
            "",
            "Next:",
            "→ How to extract features from data",
            "→ How classification works",
            "→ Practical examples"
        ],
        "Checkmark icons, summary diagram"
    )
    
    # ============================================
    # PART 2: FEATURE EXTRACTION & CLASSIFICATION
    # ============================================
    
    add_content_slide(
        "Part 2: Feature Extraction & Classification",
        [
            "Duration: 60 minutes",
            "Topics:",
            "• What is data and features?",
            "• Feature extraction methods",
            "• What is classification?",
            "• Classification examples",
            "• Common classifiers",
            "• End-to-end example"
        ],
        "Icon: Data transformation, classification symbol"
    )
    
    # Slide 16: What is Data?
    add_content_slide(
        "What is Data?",
        [
            "Data = Information we can use",
            "",
            "Types of data:",
            "",
            "1. Images",
            "   • Photos, drawings, screenshots",
            "   • Example: Cat photo, student ID card",
            "",
            "2. Text",
            "   • Emails, articles, messages",
            "   • Example: 'This is a great course!'",
            "",
            "3. Numbers",
            "   • Measurements, scores, prices",
            "   • Example: Student grades, temperature",
            "",
            "4. Mixed",
            "   • Combination of above",
            "   • Example: Student profile (name, photo, grades)"
        ],
        "Icons: Image, text document, numbers, mixed data"
    )
    
    # Slide 17: What are Features?
    add_content_slide(
        "What are Features?",
        [
            "Features = Characteristics that describe data",
            "",
            "Simple analogy:",
            "• Describing a student:",
            "  - Height (number)",
            "  - Hair color (category)",
            "  - Grade average (number)",
            "",
            "In ML:",
            "• Features help the model understand data",
            "• Good features = Better predictions",
            "",
            "Example - Image of a cat:",
            "• Raw: Millions of pixels",
            "• Features:",
            "  - Has pointy ears? (Yes/No)",
            "  - Has whiskers? (Yes/No)",
            "  - Color: Orange, Black, White",
            "  - Size: Small, Medium, Large"
        ],
        "Diagram: Raw data → Feature extraction → Features list"
    )
    
    # Slide 18: Why Feature Extraction is Important
    add_content_slide(
        "Why Feature Extraction is Important?",
        [
            "Problem with raw data:",
            "• Too much information",
            "• Irrelevant details",
            "• Hard for computer to understand",
            "",
            "Solution - Feature extraction:",
            "• Focus on important characteristics",
            "• Reduce complexity",
            "• Make patterns easier to find",
            "",
            "Example - Student photo:",
            "• Raw: 1000x1000 pixels = 1 million numbers",
            "• Features:",
            "  - Face shape: Round",
            "  - Eye color: Brown",
            "  - Hair: Long, Black",
            "  - Expression: Smiling",
            "• Result: Just 4-5 meaningful features!"
        ],
        "Diagram: Large raw data → Small feature set"
    )
    
    # Slide 19: Feature Extraction - Images
    add_content_slide(
        "Feature Extraction: Images",
        [
            "Image features we can extract:",
            "",
            "1. Edges",
            "   • Detect boundaries and shapes",
            "   • Example: Outline of a face",
            "",
            "2. Colors",
            "   • Dominant colors, color distribution",
            "   • Example: Sky is blue, grass is green",
            "",
            "3. Textures",
            "   • Surface patterns",
            "   • Example: Smooth skin, rough bark",
            "",
            "4. Shapes",
            "   • Geometric forms",
            "   • Example: Circle (ball), Rectangle (book)",
            "",
            "Real example:",
            "• Cat photo → Extract: Pointy ears, whiskers, tail"
        ],
        "Diagram: Image → Edge detection → Color analysis → Features"
    )
    
    # Slide 20: Feature Extraction - Text
    add_content_slide(
        "Feature Extraction: Text",
        [
            "Text features we can extract:",
            "",
            "1. Word Count",
            "   • How many times each word appears",
            "   • Example: 'free' appears 5 times → Spam indicator",
            "",
            "2. Word Frequency (TF-IDF - Simple explanation)",
            "   • Important words appear more often",
            "   • Example: 'machine learning' in AI article",
            "",
            "3. Text Length",
            "   • Number of characters or words",
            "   • Example: Short messages vs long essays",
            "",
            "4. Special Characters",
            "   • Exclamation marks, capital letters",
            "   • Example: 'FREE!!!' → Spam indicator",
            "",
            "Real example:",
            "• Email: 'Win FREE money!!!'",
            "• Features: Has 'free' (1), Has 'win' (1), Many '!' (3)"
        ],
        "Diagram: Text → Word analysis → Feature vector"
    )
    
    # Slide 21: Feature Extraction - Numerical Data
    add_content_slide(
        "Feature Extraction: Numerical Data",
        [
            "Numerical features we can extract:",
            "",
            "1. Statistics",
            "   • Mean (average)",
            "   • Max, Min values",
            "   • Example: Average student grade = 85",
            "",
            "2. Trends",
            "   • Increasing or decreasing over time",
            "   • Example: Temperature rising in summer",
            "",
            "3. Comparisons",
            "   • Above/below average",
            "   • Example: Grade 90 is above class average",
            "",
            "4. Categories",
            "   • Convert numbers to groups",
            "   • Example: Age 20 → 'Young adult'",
            "",
            "Real example - Student data:",
            "• Raw: [85, 90, 78, 92, 88]",
            "• Features: Mean=86.6, Max=92, Trend=Increasing"
        ],
        "Chart: Raw numbers → Statistics → Features"
    )
    
    # Slide 22: What is Classification?
    add_content_slide(
        "What is Classification?",
        [
            "Classification = Putting data into categories",
            "",
            "Simple analogy:",
            "• Sorting fruits into baskets",
            "• Apple → Fruit basket",
            "• Carrot → Vegetable basket",
            "",
            "In ML:",
            "• Input: Data with features",
            "• Output: Category/Class",
            "",
            "Example:",
            "• Input: Email with features",
            "  - Has 'free' (Yes)",
            "  - Has '!!!' (Yes)",
            "  - From unknown sender (Yes)",
            "• Output: Spam",
            "",
            "Goal:",
            "• Learn patterns from examples",
            "• Classify new data correctly"
        ],
        "Diagram: Input features → Classifier → Category output"
    )
    
    # Slide 23: Binary vs Multi-class Classification
    add_two_column_slide(
        "Binary vs Multi-class Classification",
        [
            "Binary Classification:",
            "",
            "Two categories only:",
            "• Yes/No",
            "• Spam/Not Spam",
            "• Pass/Fail",
            "• Sick/Healthy",
            "",
            "Example:",
            "• Email → Spam or Not Spam?",
            "",
            "Simpler problem:",
            "• Just two choices",
            "• Easier to understand"
        ],
        [
            "Multi-class Classification:",
            "",
            "More than two categories:",
            "• Cat/Dog/Bird",
            "• Grade: A/B/C/D/F",
            "• Emotion: Happy/Sad/Angry",
            "",
            "Example:",
            "• Image → Cat, Dog, or Bird?",
            "",
            "More complex:",
            "• Multiple choices",
            "• Need more training data"
        ],
        "Venn diagram or comparison chart"
    )
    
    # Slide 24: Classification Example - Spam Detection
    add_content_slide(
        "Classification Example: Spam Detection",
        [
            "Problem:",
            "• Classify emails as Spam or Not Spam",
            "",
            "Step 1: Extract Features",
            "• Word 'free' appears? (Yes/No)",
            "• Word 'win' appears? (Yes/No)",
            "• Many exclamation marks? (Yes/No)",
            "• From known sender? (Yes/No)",
            "",
            "Step 2: Train Classifier",
            "• Show 1000 spam emails",
            "• Show 1000 normal emails",
            "• Learn patterns",
            "",
            "Step 3: Classify New Email",
            "• Extract features",
            "• Compare with learned patterns",
            "• Output: Spam or Not Spam",
            "",
            "Result:",
            "• 99%+ accuracy in real systems"
        ],
        "Flowchart: Email → Features → Classifier → Spam/Not Spam"
    )
    
    # Slide 25: Classification Example - Disease Detection
    add_content_slide(
        "Classification Example: Disease Detection",
        [
            "Problem:",
            "• Detect disease from symptoms",
            "",
            "Features:",
            "• Fever? (Yes/No)",
            "• Cough? (Yes/No)",
            "• Temperature: 38.5°C",
            "• Age: 25 years",
            "",
            "Classification:",
            "• Input: Patient symptoms",
            "• Output: Disease category",
            "  - Common cold",
            "  - Flu",
            "  - COVID-19",
            "  - Healthy",
            "",
            "Real application:",
            "• Assist doctors in diagnosis",
            "• Early detection saves lives",
            "",
            "Note:",
            "• ML assists, doesn't replace doctors"
        ],
        "Diagram: Symptoms → Features → Classifier → Diagnosis"
    )
    
    # Slide 26: Classification Example - Face Recognition
    add_content_slide(
        "Classification Example: Face Recognition",
        [
            "Problem:",
            "• Identify person in photo",
            "",
            "Step 1: Extract Face Features",
            "• Eye distance",
            "• Nose shape",
            "• Face width/height ratio",
            "• Unique facial markers",
            "",
            "Step 2: Compare with Database",
            "• Match features with known faces",
            "• Calculate similarity score",
            "",
            "Step 3: Classify",
            "• If similarity > 95% → Match found",
            "• Output: Person's name",
            "",
            "Applications:",
            "• Phone unlock",
            "• Security systems",
            "• Photo tagging"
        ],
        "Diagram: Face → Feature extraction → Comparison → Identity"
    )
    
    # Slide 27: Common Classifiers - Overview
    add_content_slide(
        "Common Classifiers (High-Level Overview)",
        [
            "Three popular classifiers:",
            "",
            "1. K-Nearest Neighbors (KNN)",
            "   • Find similar examples",
            "   • Vote by majority",
            "   • Like asking neighbors for advice",
            "",
            "2. Decision Tree",
            "   • Ask yes/no questions",
            "   • Follow path to answer",
            "   • Like a flowchart",
            "",
            "3. Logistic Regression",
            "   • Find best line/curve to separate classes",
            "   • Uses probability",
            "   • Like drawing a boundary"
        ],
        "Three icons/diagrams for each classifier type"
    )
    
    # Slide 28: K-Nearest Neighbors (KNN)
    add_content_slide(
        "K-Nearest Neighbors (KNN) - Intuition",
        [
            "How it works (simple explanation):",
            "",
            "1. Store all training examples",
            "",
            "2. For new data:",
            "   • Find K nearest neighbors",
            "   • K = number of neighbors (e.g., K=3)",
            "",
            "3. Vote:",
            "   • Look at neighbors' labels",
            "   • Majority wins",
            "",
            "Example:",
            "• New email: Is it spam?",
            "• Find 3 most similar emails",
            "• 2 are spam, 1 is not",
            "• Decision: Spam (2 out of 3)",
            "",
            "Analogy:",
            "• Like asking your 3 closest friends for advice"
        ],
        "Diagram: New point → Find K neighbors → Vote → Class"
    )
    
    # Slide 29: Decision Tree
    add_content_slide(
        "Decision Tree",
        [
            "How it works:",
            "",
            "• Series of yes/no questions",
            "• Each answer leads to next question",
            "• Final answer at the end",
            "",
            "Example - Classifying email:",
            "",
            "Question 1: Has word 'free'?",
            "  → Yes: Go to Question 2",
            "  → No: Go to Question 3",
            "",
            "Question 2: From unknown sender?",
            "  → Yes: SPAM",
            "  → No: Go to Question 3",
            "",
            "Question 3: Has many '!'?",
            "  → Yes: SPAM",
            "  → No: NOT SPAM",
            "",
            "Visual:",
            "• Like a tree with branches",
            "• Each branch = a decision"
        ],
        "Tree diagram: Root → Branches → Leaves (classes)"
    )
    
    # Slide 30: Logistic Regression (Concept)
    add_content_slide(
        "Logistic Regression (Concept Only)",
        [
            "High-level idea:",
            "",
            "• Draw a line/curve to separate classes",
            "• Find the best boundary",
            "",
            "Example - Spam detection:",
            "• Feature 1: Number of 'free' words",
            "• Feature 2: Number of exclamation marks",
            "• Draw line separating spam from not spam",
            "",
            "How it works:",
            "• Calculate probability",
            "• If probability > 50% → Class 1",
            "• If probability < 50% → Class 2",
            "",
            "Visual:",
            "• Like drawing a line on a graph",
            "• Points on one side = Spam",
            "• Points on other side = Not Spam",
            "",
            "Note:",
            "• We'll skip the math details",
            "• Focus on the concept"
        ],
        "Graph: Two classes separated by a curve/line"
    )
    
    # Slide 31: End-to-End Example
    add_content_slide(
        "End-to-End Example: Image Classification",
        [
            "Complete workflow:",
            "",
            "Step 1: Input Image",
            "• Photo of a cat",
            "",
            "Step 2: Feature Extraction",
            "• Extract: Pointy ears (Yes)",
            "• Extract: Whiskers (Yes)",
            "• Extract: Tail (Yes)",
            "• Extract: Size (Small)",
            "",
            "Step 3: Classification",
            "• Input features to classifier",
            "• Compare with learned patterns",
            "",
            "Step 4: Output",
            "• Prediction: Cat (95% confidence)",
            "",
            "Real system:",
            "• Processes in milliseconds",
            "• Can handle thousands of images"
        ],
        "Complete flowchart: Image → Features → Classifier → Output"
    )
    
    # Slide 32: Student Example - Grade Prediction
    add_content_slide(
        "Student Example: Grade Prediction",
        [
            "Problem:",
            "• Predict if student will pass or fail",
            "",
            "Features:",
            "• Attendance: 85%",
            "• Assignment average: 78",
            "• Midterm score: 82",
            "• Study hours per week: 15",
            "",
            "Classification:",
            "• Input: Student features",
            "• Output: Pass or Fail",
            "",
            "How it helps:",
            "• Identify at-risk students early",
            "• Provide extra support",
            "• Improve learning outcomes",
            "",
            "Real application:",
            "• Used in many universities",
            "• Helps improve student success"
        ],
        "Diagram: Student data → Features → Classifier → Pass/Fail"
    )
    
    # Slide 33: Part 2 Recap
    add_content_slide(
        "Part 2: Key Takeaways",
        [
            "What we learned:",
            "",
            "✓ Features = Important characteristics of data",
            "✓ Feature extraction simplifies complex data",
            "✓ Classification = Putting data into categories",
            "✓ Binary (2 classes) vs Multi-class (many classes)",
            "✓ Common classifiers: KNN, Decision Tree, Logistic Regression",
            "",
            "Complete process:",
            "Data → Feature Extraction → Classification → Prediction",
            "",
            "Next steps:",
            "→ Practice with real datasets",
            "→ Try building your own classifier",
            "→ Explore more ML topics"
        ],
        "Summary diagram, checkmark icons"
    )
    
    # Slide 34: Final Summary
    add_content_slide(
        "Course Summary",
        [
            "Today we covered:",
            "",
            "Part 1: Introduction to ML",
            "• What is ML and why we need it",
            "• Types: Supervised, Unsupervised, Reinforcement",
            "• Real-world applications",
            "• ML workflow",
            "",
            "Part 2: Feature Extraction & Classification",
            "• Understanding data and features",
            "• Feature extraction methods",
            "• Classification concepts",
            "• Common classifiers",
            "",
            "Key message:",
            "ML is about learning patterns from data to make predictions"
        ],
        "Summary diagram connecting all concepts"
    )
    
    # Slide 35: Questions & Next Steps
    add_content_slide(
        "Questions & Next Steps",
        [
            "Questions?",
            "",
            "Recommended next topics:",
            "• Model evaluation (accuracy, precision, recall)",
            "• More classifiers (SVM, Random Forest)",
            "• Deep Learning basics",
            "• Practical projects",
            "",
            "Resources:",
            "• Practice with datasets (Kaggle, UCI)",
            "• Try Python libraries (scikit-learn)",
            "• Build simple projects",
            "",
            "Thank you!",
            "",
            "Remember:",
            "• Start simple",
            "• Practice regularly",
            "• Learn by doing"
        ],
        "Question mark icon, resource icons"
    )
    
    return prs

# Generate and save the presentation
if __name__ == "__main__":
    print("Creating Machine Learning lecture presentation...")
    presentation = create_ml_presentation()
    output_file = "Artificial Intelligence_Husain/Notes/ML_Lecture_Introduction_and_Features.pptx"
    presentation.save(output_file)
    print(f"Presentation saved to: {output_file}")
    print(f"Total slides: {len(presentation.slides)}")

